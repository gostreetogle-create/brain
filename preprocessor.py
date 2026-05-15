import os
import hashlib
from config import Config

try:
    import pytesseract
    from PIL import Image
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


def file_hash(filepath: str) -> str:
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".bmp"}


def extract_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    size = os.path.getsize(filepath)
    if size > Config.MAX_FILE_SIZE_BYTES:
        raise ValueError(f"File too large: {size} bytes (max {Config.MAX_FILE_SIZE_BYTES})")

    if ext in IMAGE_EXTS:
        return _extract_image(filepath)

    extractors = {
        ".pdf": _extract_pdf,
        ".txt": _extract_text,
        ".md": _extract_text,
        ".csv": _extract_text,
        ".json": _extract_text,
        ".xml": _extract_text,
        ".html": _extract_text,
        ".htm": _extract_text,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
        ".xls": _extract_xls,
    }

    extractor = extractors.get(ext, _extract_fallback)
    return extractor(filepath)


def _extract_text(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_pdf(filepath: str) -> str:
    if not PYMUPDF_AVAILABLE:
        return _extract_fallback(filepath)
    text = []
    with fitz.open(filepath) as doc:
        for page in doc:
            text.append(page.get_text())
    result = "\n".join(text).strip()
    if result:
        return result
    if TESSERACT_AVAILABLE:
        return _ocr_pdf(filepath)
    return result


def _ocr_pdf(filepath: str) -> str:
    import fitz
    text = []
    with fitz.open(filepath) as doc:
        for page in doc:
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text.append(pytesseract.image_to_string(img, lang="rus+eng"))
    return "\n".join(text)


def _extract_docx(filepath: str) -> str:
    if not DOCX_AVAILABLE:
        return _extract_fallback(filepath)
    doc = Document(filepath)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(filepath: str) -> str:
    if not PPTX_AVAILABLE:
        return _extract_fallback(filepath)
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def _extract_xlsx(filepath: str) -> str:
    if not EXCEL_AVAILABLE:
        return _extract_fallback(filepath)
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            texts.append("\t".join(str(c) for c in row if c is not None))
    return "\n".join(texts)


def _extract_xls(filepath: str) -> str:
    if not XLRD_AVAILABLE:
        return _extract_fallback(filepath)
    wb = xlrd.open_workbook(filepath)
    texts = []
    for sheet_idx in range(wb.nsheets):
        ws = wb.sheet_by_index(sheet_idx)
        for row_idx in range(ws.nrows):
            texts.append("\t".join(str(ws.cell_value(row_idx, c)) for c in range(ws.ncols)))
    return "\n".join(texts)


def _tesseract_installed() -> bool:
    if not TESSERACT_AVAILABLE:
        return False
    try:
        import subprocess
        subprocess.run(["tesseract", "--version"], capture_output=True, timeout=5)
        return True
    except Exception:
        return False


def _extract_image(filepath: str) -> str:
    if not _tesseract_installed():
        raise ValueError(
            "Для обработки изображений нужен Tesseract OCR.\n"
            "Установи: https://github.com/UB-Mannheim/tesseract/wiki\n"
            "Или пересохрани файл в PDF/текст."
        )
    img = Image.open(filepath)
    return pytesseract.image_to_string(img, lang="rus+eng")


def _extract_fallback(filepath: str) -> str:
    try:
        return _extract_text(filepath)
    except Exception:
        return f"[Не удалось извлечь текст: {os.path.basename(filepath)}]"
